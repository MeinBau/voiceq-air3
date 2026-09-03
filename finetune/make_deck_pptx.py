"""발표자료 PPTX 생성기 — HTML 슬라이드와 같은 내용을 PowerPoint 파일로 낸다.

HTML 쪽(artifact)이 발표용이고 이 파일은 제출·편집용이다. PowerPoint에서 열어
문구를 고치거나 팀원이 슬라이드를 더 넣을 수 있도록, 텍스트는 전부 편집 가능한
도형으로 넣는다(이미지로 굽지 않는다).

수치는 전부 실측이다. 출처는 finetune/README.md 4절.

사용:
    python finetune/make_deck_pptx.py                 # ./voicecue_발표.pptx
    python finetune/make_deck_pptx.py --out 경로.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# HTML 덱과 같은 팔레트 — 어두운 지휘소 표시장치 톤에 호박색 강조 하나.
GROUND = RGBColor(0x0A, 0x10, 0x14)
PANEL = RGBColor(0x12, 0x1B, 0x21)
LINE = RGBColor(0x24, 0x34, 0x3D)
INK = RGBColor(0xE6, 0xEE, 0xF2)
INK_DIM = RGBColor(0x93, 0xA8, 0xB3)
INK_FAINT = RGBColor(0x61, 0x75, 0x7F)
SIGNAL = RGBColor(0xE0, 0xA9, 0x3B)

# 16:9
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(760000)

HEAD_FONT = "맑은 고딕"
BODY_FONT = "맑은 고딕"
MONO_FONT = "Consolas"


def cm(v: float) -> Emu:
    return Emu(int(v * 360000))


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.n = 0

    # ---------- 기본 골격 ----------
    def _blank(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = GROUND
        return s

    def _text(self, slide, x, y, w, h, runs, *, size=14, color=INK, bold=False,
              font=BODY_FONT, align=PP_ALIGN.LEFT, spacing=1.35, anchor=MSO_ANCHOR.TOP):
        """runs: 문자열 또는 (텍스트, {속성}) 목록."""
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        items = runs if isinstance(runs, list) else [runs]
        first = True
        for item in items:
            text, opts = item if isinstance(item, tuple) else (item, {})
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = opts.get("align", align)
            p.line_spacing = opts.get("spacing", spacing)
            if opts.get("space_before"):
                p.space_before = Pt(opts["space_before"])
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", font)
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.color.rgb = opts.get("color", color)
        return box

    def _card(self, slide, x, y, w, h, *, accent=False):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = PANEL
        shp.line.color.rgb = SIGNAL if accent else LINE
        shp.line.width = Pt(1)
        shp.shadow.inherit = False
        return shp

    def _chrome(self, slide, section: str):
        self.n += 1
        self._text(slide, MARGIN, cm(0.62), cm(18), cm(0.6),
                   [("VOICE-CUE", {"color": SIGNAL, "bold": True}),
                    (f"   {section}", {"color": INK_FAINT})],
                   size=10, font=MONO_FONT, spacing=1.0)
        self._text(slide, SLIDE_W - MARGIN - cm(6), SLIDE_H - cm(1.5), cm(6), cm(0.6),
                   f"{self.n:02d}", size=10, color=INK_FAINT, font=MONO_FONT,
                   align=PP_ALIGN.RIGHT, spacing=1.0)

    # ---------- 슬라이드 유형 ----------
    def title_slide(self, eyebrow, title, lede, figures):
        s = self._blank()
        self._chrome(s, "표지")
        y = cm(4.0)
        self._text(s, MARGIN, y, SLIDE_W - MARGIN * 2, cm(1), eyebrow,
                   size=11, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
        self._text(s, MARGIN, y + cm(1.1), SLIDE_W - MARGIN * 2, cm(3.4), title,
                   size=40, color=INK, bold=True, font=HEAD_FONT, spacing=1.15)
        self._text(s, MARGIN, y + cm(4.9), cm(22), cm(2), lede,
                   size=15, color=INK_DIM, spacing=1.5)
        self._figure_row(s, MARGIN, cm(12.6), figures)
        return s

    def _figure_row(self, slide, x, y, figures, width=None):
        colw = (width or (SLIDE_W - MARGIN * 2)) // len(figures)
        for i, (v, k) in enumerate(figures):
            cx = x + colw * i
            self._text(slide, cx, y, colw - cm(0.4), cm(1.5), v,
                       size=30, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
            self._text(slide, cx, y + cm(1.5), colw - cm(0.4), cm(1.2), k,
                       size=10, color=INK_FAINT, spacing=1.3)

    def content(self, section, eyebrow, heading, blocks, *, figures=None, note=None):
        """blocks: (제목, 본문) 카드 목록. 2~4개면 가로 배치."""
        s = self._blank()
        self._chrome(s, section)
        self._text(s, MARGIN, cm(2.1), SLIDE_W - MARGIN * 2, cm(0.7), eyebrow,
                   size=10, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
        self._text(s, MARGIN, cm(2.9), SLIDE_W - MARGIN * 2, cm(2.2), heading,
                   size=26, color=INK, bold=True, font=HEAD_FONT, spacing=1.2)

        top = cm(6.0)
        avail_h = cm(9.4) if not (figures or note) else cm(7.4)
        if blocks:
            n = len(blocks)
            gap = cm(0.5)
            cw = (SLIDE_W - MARGIN * 2 - gap * (n - 1)) // n
            for i, (t, body) in enumerate(blocks):
                cx = MARGIN + (cw + gap) * i
                self._card(s, cx, top, cw, avail_h)
                self._text(s, cx + cm(0.6), top + cm(0.6), cw - cm(1.2), cm(1.4), t,
                           size=14, color=INK, bold=True, spacing=1.3)
                self._text(s, cx + cm(0.6), top + cm(1.9), cw - cm(1.2), avail_h - cm(2.5),
                           body, size=11, color=INK_DIM, spacing=1.45)

        if figures:
            self._figure_row(s, MARGIN, top + avail_h + cm(0.7), figures)
        if note:
            self._text(s, MARGIN, SLIDE_H - cm(3.4), SLIDE_W - MARGIN * 2, cm(2.2), note,
                       size=10, color=INK_FAINT, spacing=1.4)
        return s

    def table_slide(self, section, eyebrow, heading, header, rows, *, note=None,
                    col_widths=None, highlight_col=None):
        s = self._blank()
        self._chrome(s, section)
        self._text(s, MARGIN, cm(2.1), SLIDE_W - MARGIN * 2, cm(0.7), eyebrow,
                   size=10, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
        self._text(s, MARGIN, cm(2.9), SLIDE_W - MARGIN * 2, cm(2.0), heading,
                   size=26, color=INK, bold=True, font=HEAD_FONT, spacing=1.2)

        n_rows, n_cols = len(rows) + 1, len(header)
        tw = SLIDE_W - MARGIN * 2
        th = cm(0.85) * n_rows
        gt = s.shapes.add_table(n_rows, n_cols, MARGIN, cm(6.0), tw, th).table
        if col_widths:
            for i, frac in enumerate(col_widths):
                gt.columns[i].width = Emu(int(tw * frac))

        for c, name in enumerate(header):
            cell = gt.cell(0, c)
            cell.text = name
            cell.fill.solid()
            cell.fill.fore_color.rgb = GROUND
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = MONO_FONT
            p.font.color.rgb = INK_FAINT
            p.font.bold = True

        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = gt.cell(r, c)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL if r % 2 else GROUND
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.name = BODY_FONT
                hot = highlight_col is not None and c == highlight_col
                p.font.color.rgb = SIGNAL if hot else INK_DIM
                p.font.bold = bool(hot)

        if note:
            self._text(s, MARGIN, SLIDE_H - cm(3.6), SLIDE_W - MARGIN * 2, cm(2.4), note,
                       size=10, color=INK_FAINT, spacing=1.4)
        return s

    def metric_slide(self, section, eyebrow, heading, metrics, *, note=None):
        """튜닝 전/후 막대 비교."""
        from pptx.enum.shapes import MSO_SHAPE
        s = self._blank()
        self._chrome(s, section)
        self._text(s, MARGIN, cm(2.1), SLIDE_W - MARGIN * 2, cm(0.7), eyebrow,
                   size=10, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
        self._text(s, MARGIN, cm(2.9), SLIDE_W - MARGIN * 2, cm(2.0), heading,
                   size=26, color=INK, bold=True, font=HEAD_FONT, spacing=1.2)

        y = cm(6.2)
        label_w, track_w, row_h = cm(6.4), cm(15.5), cm(1.25)
        for label, before, after in metrics:
            self._text(s, MARGIN, y + cm(0.12), label_w, cm(0.8), label,
                       size=12, color=INK_DIM, spacing=1.0)
            bx = MARGIN + label_w
            base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, y, track_w, cm(0.62))
            base.fill.solid(); base.fill.fore_color.rgb = PANEL
            base.line.fill.background(); base.shadow.inherit = False

            bw = int(track_w * before / 100)
            b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, y, bw, cm(0.62))
            b.fill.solid(); b.fill.fore_color.rgb = LINE
            b.line.fill.background(); b.shadow.inherit = False

            aw = int(track_w * after / 100)
            a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, y + cm(0.16), aw, cm(0.3))
            a.fill.solid(); a.fill.fore_color.rgb = SIGNAL
            a.line.fill.background(); a.shadow.inherit = False

            self._text(s, bx + track_w + cm(0.35), y + cm(0.1), cm(4), cm(0.7),
                       f"{before:.1f} → {after:.1f}  (+{after - before:.1f})",
                       size=11, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
            y += row_h

        self._text(s, MARGIN + label_w, y + cm(0.3), cm(16), cm(0.6),
                   "밝은 막대 = 튜닝 후 · 어두운 막대 = 튜닝 전",
                   size=10, color=INK_FAINT, font=MONO_FONT, spacing=1.0)
        if note:
            self._text(s, MARGIN, SLIDE_H - cm(3.2), SLIDE_W - MARGIN * 2, cm(2.2), note,
                       size=10, color=INK_FAINT, spacing=1.4)
        return s

    def flow_slide(self, section, eyebrow, heading, nodes, *, blocks=None, note=None):
        s = self._blank()
        self._chrome(s, section)
        self._text(s, MARGIN, cm(2.1), SLIDE_W - MARGIN * 2, cm(0.7), eyebrow,
                   size=10, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
        self._text(s, MARGIN, cm(2.9), SLIDE_W - MARGIN * 2, cm(2.0), heading,
                   size=26, color=INK, bold=True, font=HEAD_FONT, spacing=1.2)

        n = len(nodes)
        gap = cm(0.75)
        cw = (SLIDE_W - MARGIN * 2 - gap * (n - 1)) // n
        top, ch = cm(6.0), cm(4.2)
        for i, (t, d, accent) in enumerate(nodes):
            cx = MARGIN + (cw + gap) * i
            self._card(s, cx, top, cw, ch, accent=accent)
            self._text(s, cx + cm(0.5), top + cm(0.5), cw - cm(1.0), cm(1.2), t,
                       size=13, color=SIGNAL if accent else INK, bold=True, spacing=1.25)
            self._text(s, cx + cm(0.5), top + cm(1.7), cw - cm(1.0), ch - cm(2.2), d,
                       size=10, color=INK_FAINT, spacing=1.4)
            if i < n - 1:
                self._text(s, cx + cw, top + cm(1.7), gap, cm(0.8), "→",
                           size=14, color=INK_FAINT, font=MONO_FONT,
                           align=PP_ALIGN.CENTER, spacing=1.0)

        if blocks:
            by, bh = top + ch + cm(0.8), cm(3.6)
            m = len(blocks)
            bw = (SLIDE_W - MARGIN * 2 - gap * (m - 1)) // m
            for i, (t, body) in enumerate(blocks):
                cx = MARGIN + (bw + gap) * i
                self._card(s, cx, by, bw, bh)
                self._text(s, cx + cm(0.5), by + cm(0.45), bw - cm(1.0), cm(1.0), t,
                           size=12, color=INK, bold=True, spacing=1.25)
                self._text(s, cx + cm(0.5), by + cm(1.4), bw - cm(1.0), bh - cm(1.9), body,
                           size=10, color=INK_DIM, spacing=1.4)
        if note:
            self._text(s, MARGIN, SLIDE_H - cm(3.2), SLIDE_W - MARGIN * 2, cm(2.2), note,
                       size=10, color=INK_FAINT, spacing=1.4)
        return s

    def toc(self, items):
        """목차 — (번호, 제목, 설명, 쪽수) 목록을 2열로."""
        s = self._blank()
        self._chrome(s, "목차")
        self._text(s, MARGIN, cm(2.1), SLIDE_W - MARGIN * 2, cm(0.7), "목차",
                   size=10, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
        self._text(s, MARGIN, cm(2.9), SLIDE_W - MARGIN * 2, cm(1.6), "오늘 말씀드릴 순서",
                   size=26, color=INK, bold=True, font=HEAD_FONT, spacing=1.2)

        colw = (SLIDE_W - MARGIN * 2 - cm(1.2)) // 2
        row_h = cm(2.15)
        for i, (num, title, desc, pages) in enumerate(items):
            cx = MARGIN + (colw + cm(1.2)) * (i % 2)
            cy = cm(5.4) + row_h * (i // 2)
            self._text(s, cx, cy, cm(1.4), cm(0.8), num,
                       size=13, color=SIGNAL, bold=True, font=MONO_FONT, spacing=1.0)
            self._text(s, cx + cm(1.5), cy, colw - cm(1.5), cm(0.8),
                       [(title, {"size": 14, "color": INK, "bold": True}),
                        (f"   {pages}", {"size": 10, "color": INK_FAINT, "font": MONO_FONT})],
                       spacing=1.0)
            self._text(s, cx + cm(1.5), cy + cm(0.75), colw - cm(1.5), cm(1.2), desc,
                       size=10, color=INK_DIM, spacing=1.35)
        return s

    def closing(self, heading, figures, lede):
        s = self._blank()
        self._chrome(s, "마무리")
        self._text(s, MARGIN, cm(3.6), SLIDE_W - MARGIN * 2, cm(3.0), heading,
                   size=34, color=INK, bold=True, font=HEAD_FONT, spacing=1.2)
        self._figure_row(s, MARGIN, cm(8.2), figures, width=cm(24))
        self._text(s, MARGIN, cm(11.8), SLIDE_W - MARGIN * 2, cm(3.4), lede,
                   size=13, color=INK_DIM, spacing=1.5)
        return s

    def save(self, path: Path) -> None:
        self.prs.save(str(path))


# 실측값 — Qwen2.5-1.5B-Instruct, 1에폭, test 40턴, 튜닝 전후 모두 4bit 로드
METRICS = [
    ("상황유형 정확도", 45.00, 75.00),
    ("COP 셀 일치율", 73.75, 88.33),
    ("일지 키워드 정확도", 60.59, 84.93),
    ("일지 분류 정확도", 82.50, 92.50),
    ("ROUGE-L", 54.90, 79.21),
]


def build() -> Deck:
    d = Deck()

    d.title_slide(
        "제8회 공군 창의·혁신 아이디어 공모 해커톤 · 항공우주작전",
        "전투지휘소가 스스로 화면을 구성한다",
        "VOICE-CUE — 회의 발언을 알아듣고 공통작전상황도(COP)를 스스로 구성하며 "
        "작전상황일지를 대신 쓰는 AI 보조 참모.  팀 작비스(작전과 자비스)",
        [("1~2분", "현행 화면 표출 소요"), ("5초", "VOICE-CUE 목표"),
         ("1,527", "구축한 학습 데이터(턴)"), ("2GB", "폐쇄망 구동 모델 크기")],
    )

    d.toc([
        ("01", "문제 — 전투지휘소의 1~2분", "현장에서 목격한 공백 · 왜 개인의 실수가 아닌가 · 통합과 표출의 간극", "04–08"),
        ("02", "해결 구조 — 설계 판단 넷", "표출/기록 분리 · 플레이북 · 사태 단위 기록 · 화자 영향력", "09–13"),
        ("03", "시연 — 실제로 동작하는 것", "COP 자동 구성 · 작전상황일지 · 전장상황도 · 수동 보정", "14–17"),
        ("04", "파인튜닝 — 폐쇄망 sLLM", "왜 필요한가 · 데이터 구축 · 학습 구성", "18–20"),
        ("05", "성과 — 실측 결과와 품질지표", "튜닝 전/후 비교 · 기획서 4개 지표 대비 · 폐쇄망 전환", "21–23"),
        ("06", "보안 — 군 환경 설계", "3중 방어 구조 · 위험요소와 대비책", "24–25"),
        ("07", "마무리 — 기대 효과와 확장성", "작전적·인력 효과 · 타군·타 기관 수평 확산", "26–27"),
    ])

    d.content(
        "서론", "한 줄 요약",
        "만들어진 정보를 무엇을·어디에·언제까지 보여줄지 결정하는 계층",
        [("기능 ① 음성 인식",
          "발언 내용을 추출하고 발화자를 구분해 회의록으로 정리합니다. "
          "화자의 계급·직책이 이후 판단의 가중치가 됩니다."),
         ("기능 ② 작전상황일지 · 상황판",
          "회의록을 주요사태(MSEL)별로 묶어 공식 기록부를 쓰고, 사태별 중요도를 평가해 "
          "우선순위 현황판을 운영합니다."),
         ("기능 ③ 공통작전상황도 구성",
          "중요도·긴급도에 따라 정보를 자동 선택하고 Video Wall의 화면 크기·위치를 "
          "스스로 배치합니다.")],
    )

    d.content(
        "문제", "선정 배경 · 현장에서 목격한 것",
        "“무인기 2대 식별” 보고 이후, 지휘관은 1~2분간 말로만 판단한다",
        [("현행 절차 — 네 단계의 수작업",
          "① 작전병이 CCTV 프로그램을 연다\n② 해당 카메라를 찾는다\n"
          "③ 화면 크기와 위치를 조정한다\n④ 작전상황일지에 수기 입력한다"),
         ("그동안 벌어지는 일",
          "지휘관이 판단에 필요한 것은 지금 그 무인기가 보이는 화면입니다. "
          "그러나 그 화면이 뜰 때까지 눈으로 확인하지 못한 채 말로 들은 정보만으로 "
          "판단해야 하고, 화면 구성의 품질은 담당 인원의 숙련도에 따라 크게 달라집니다.")],
        figures=[("1~2분", "보고 → 화면 표출까지"), ("담당자별 상이", "화면 구성 품질")],
    )

    d.table_slide(
        "문제", "문제 · 그 1~2분 동안 벌어지는 일", "상황은 기다려주지 않는다",
        ["경과", "무슨 일이 일어나는가", "왜 문제인가"],
        [["T+0초", "“무인기 2대, 지상 차량 3대 식별되었습니다”", "지휘관이 판단을 시작해야 하는 시점"],
         ["T+10초", "작전병이 CCTV 프로그램을 연다", "어느 카메라인지는 아직 모른다"],
         ["T+30초", "북서방을 보는 카메라를 목록에서 찾는다", "카메라가 수백 대다"],
         ["T+60초", "화면 크기와 위치를 조정한다", "동시에 다른 보고가 계속 들어온다"],
         ["T+90초", "작전상황일지에 수기로 입력한다", "화면 조작과 동시에 할 수 없다"],
         ["T+120초", "이제야 화면이 뜬다", "그동안 상황은 이미 다음 단계로 넘어갔다"]],
        col_widths=[0.13, 0.47, 0.40], highlight_col=0,
        note="지휘관은 이 2분 내내 눈으로 확인하지 못한 채 말로 들은 정보만으로 판단합니다. "
             "화면이 떴을 때는 이미 다음 보고가 들어와 있고, 방금 띄운 화면은 더 이상 지금 필요한 화면이 아닐 수 있습니다.",
    )

    d.content(
        "문제", "문제 · 담당자를 탓할 수 없는 이유",
        "더 숙련된 사람을 앉혀도 풀리지 않는다",
        [("한계 ① 한 사람이 동시에 할 수 없다",
          "상황을 듣고, 카메라를 찾고, 화면을 배치하고, 일지를 쓰는 네 가지가 같은 시각에 "
          "필요합니다. 사람은 순서대로 할 수밖에 없고, 그 순서가 곧 지연입니다."),
         ("한계 ② 고를 것이 계속 늘어난다",
          "「지능형 스마트 비행단」으로 센서·영상·데이터가 과거 대비 현저히 증가했습니다. "
          "정보가 늘수록 고르는 일은 쉬워지지 않고 어려워집니다. 통합의 성과가 표출의 부담으로 "
          "되돌아옵니다."),
         ("한계 ③ 지식이 사람에게만 있다",
          "“이 상황엔 이 화면”이라는 판단은 경험에 축적된 암묵지입니다. 담당자가 바뀌면 품질도 "
          "바뀌고, 그 지식은 부대에 남지 않습니다.")],
        note="세 가지 모두 개인의 역량 문제가 아니라 구조의 문제입니다. "
             "그래서 교육이나 인력 충원이 아니라 그 판단을 체계에 옮기는 것으로 풀어야 합니다.",
    )

    d.content(
        "문제", "당면 과제",
        "정보의 ‘통합’은 진전되었으나 ‘표출 방식’은 미정립 상태다",
        [("이미 해결된 것",
          "「지능형 스마트 비행단」 구축으로 지휘관이 활용 가능한 센서·영상·데이터가 "
          "과거 대비 현저히 증가했습니다. 정보의 통합·시현 기반은 마련되었습니다."),
         ("남아 있는 것",
          "10면 이상 다중 화면에 무엇을·어디에·어떤 규모로 띄울지는 여전히 수작업입니다. "
          "작전병·정통대대 인원의 개인 역량에 좌우되고, 통합 정보량이 늘수록 화면 구성의 "
          "중요성과 난이도는 더 올라갑니다.")],
        note="정보를 더 모으는 것으로는 풀리지 않는 문제입니다. "
             "모은 정보 중에서 지금 필요한 것을 고르는 단계가 비어 있습니다.",
    )

    d.content(
        "문제", "차별화 · 기존 체계와 무엇이 다른가",
        "연주자가 아니라 지휘자다",
        [("기존 체계 — 정보를 만든다",
          "CCTV 영상 자동 감시 보고(영상→텍스트), AI 항공통제 등. "
          "각자 자기 정보를 생산하는 연주자에 해당합니다."),
         ("VOICE-CUE — 만들어진 정보를 배치한다",
          "다수의 정보를 발언·맥락에 따라 무엇을, 어느 규모로, 언제까지 보여줄지 결정하는 "
          "시현 오케스트레이션(Display Orchestration) 계층입니다. ‘지휘결심 가시화 체계’ "
          "상위에 부가되는 지능형 표출·기록 계층으로, 기존작과 중복되지 않는 신규 영역입니다.")],
    )

    d.flow_slide(
        "설계", "전체 데이터 흐름", "발언 하나가 네 가지 산출물이 되기까지",
        [("① 발언 · 화자", "편제 드롭다운 또는 Whisper 음성 인식. 계급·담당분야·영향력이 함께 실린다", False),
         ("② AI 판단", "이전 Context Memory와 함께 입력받아 주요사태를 구분하고 중요도·긴급도를 판단", True),
         ("③ Context Memory 먼저 갱신", "판단 결과와 상황 요약을 저장. 이 메모를 바탕으로 나머지가 산출된다", False),
         ("④ 세 가지 산출", "화면 구성(COP) · 작전상황판 · 작전상황일지", False)],
        blocks=[("순서가 핵심입니다",
                 "Context Memory를 가장 먼저 갱신하고 그것을 바탕으로 나머지를 산출하므로 앞선 상황이 "
                 "일관되게 반영됩니다. 갱신된 메모는 다음 발언에서 다시 입력으로 읽힙니다."),
                ("사람이 보는 기록과 AI의 기억은 별개입니다",
                 "작전상황일지는 사람이 읽는 공식 기록부이고 Context Memory는 AI가 다음 판단에 쓰는 "
                 "내부 기억입니다. 둘을 분리해야 일지가 AI의 중간 사고로 오염되지 않습니다.")],
    )

    d.content(
        "설계", "핵심 설계 ① · 표출 지연",
        "화면에 필요한 것만 따로 먼저 받는다",
        [("FAST — 표출 경로",
          "상황 유형 하나만 받습니다. 출력이 짧아 빠르고, 화면은 이것이 도착하는 즉시 갱신됩니다.\n\n"
          '{ "situation": { "type": "드론상황" } }'),
         ("FULL — 기록 경로",
          "Context Memory · 상황판 · 작전상황일지를 만듭니다. 출력이 길어 느리지만 "
          "화면 표출을 붙잡지 않습니다."),
         ("왜 나눴나",
          "한 번의 호출로 네 가지를 다 만들면 가장 느린 산출물이 화면 표출까지 붙잡아 둡니다. "
          "두 호출을 스레드로 동시에 던지면 느린 쪽 하나만큼만 걸립니다.")],
        note="기획서 3-라① “필요정보 표출 지연시간 5초 이내”는 FAST 경로의 지연으로 측정합니다.",
    )

    d.flow_slide(
        "설계", "핵심 설계 ② · 배치 품질", "LLM은 화면을 고르지 않는다",
        [("LLM", "상황 유형 분류만 한다 (11종 중 하나)", False),
         ("운용자 플레이북", "상황 유형별 화면 슬롯 정답 레이아웃. 앱에서 표로 편집한다", True),
         ("코드가 슬롯 해석", "발언의 방위·시설명과 태그가 겹치는 CCTV를 직접 검색", False),
         ("폐쇄 카탈로그", "화면 소스 270건 중에서만 선택. 없는 화면을 지어낼 수 없다", False)],
        blocks=[("왜 이렇게 했나",
                 "배치를 LLM이 매번 자유롭게 정하면 같은 상황에도 결과가 흔들립니다. 운용 경험자가 "
                 "“이 상황엔 이 화면을 이 순서로”를 이미 알고 있으므로 그 판단을 데이터로 고정했습니다."),
                ("지표 ②가 구조적으로 보장됩니다",
                 "정답이 곧 플레이북이므로, 플레이북을 따르는 한 전문가 정답 레이아웃 대비 일치율은 "
                 "정의상 100%입니다. 남는 평가 대상은 “상황 유형을 맞게 분류했는가” 하나로 좁혀집니다.")],
    )

    d.table_slide(
        "설계", "핵심 설계 ③ · 기록 품질", "발언 하나가 곧 사태 하나가 아니다",
        ["분류", "발언", "일지 반영"],
        [["상황", "북서방 상공에 무인기 2대 식별되었습니다", "사태1 신설"],
         ["조치", "북측 대공포 사격대기 완료했습니다", "사태1에 병합"],
         ["조치", "1번은 격추 대응, 2번은 추적 유지합니다", "사태1에 병합"],
         ["조치", "차량 3대는 정상 지원 차량으로 확인되었습니다", "사태1에 병합(정정)"],
         ["무시", "저건 왜 저렇게 대응하는 겁니까?", "일지 불변"]],
        col_widths=[0.10, 0.58, 0.32], highlight_col=2,
        note="하나의 사태는 여러 발언에 걸쳐 진행됩니다. 발언마다 새 항목을 만들면 일지가 "
             "사태가 아니라 발언의 나열이 되어 버립니다.  "
             "학습 데이터도 이 비율을 반영합니다 — 상황 242 · 조치 838 · 무시 135건.",
    )

    d.table_slide(
        "설계", "핵심 설계 ④ · 발화자 구분", "누가 말했는지가 판단의 무게를 바꾼다",
        ["직책", "담당분야", "영향력", "취급"],
        [["비행단장 (준장)", "지휘 · 결심 · 교전승인", "1.00", "결심으로 즉시 최우선 반영"],
         ["기지방호전대장 (대령)", "방호 · 경계 · 대공방어", "0.85", "결심"],
         ["정보과장 (중령)", "정보 · 위협판단", "0.75", "담당분야 내 권위 있는 판단"],
         ["대공방어대장 (소령)", "대공방어 · 교전", "0.65", "담당분야 내 권위 있는 판단"],
         ["지휘통제망 채팅", "상황전파 · 신규상황", "0.60", "신규 사태 개시 근거"],
         ["기타 지휘관", "일반 · 회의참석", "0.50", "사실 전파 — 단독으로 판단을 뒤집지 않음"]],
        col_widths=[0.24, 0.26, 0.11, 0.39], highlight_col=2,
        note="화자 이름만 넘기면 LLM은 그가 누구인지 모릅니다. 편제 20개 직책·21개 부대를 정의해 "
             "프롬프트에 그대로 주입합니다.  "
             "신규 사태는 실제 운용대로 항작 계열 또는 지휘통제망 채팅이 전파합니다(데이터의 83%).",
    )

    d.content(
        "시연", "시연 ① · 공통작전상황도 자동 구성",
        "“북서방 상공에 무인기 2대 식별” → 즉시 이 배치",
        [("Video Wall 2행 6열",
          "1순위 기지 전장상황도 — 4칸 (고정)\n2순위 북서 열상감시장비(TOD) — 2칸\n"
          "3순위 SSR 화면 — 2칸\n4순위 비행 스케줄 화면 — 2칸\n"
          "5순위 작전상황판 — 1칸 (상시)\n6순위 대공방어 자산 현황 — 1칸 (상시)"),
         ("배치 규칙",
          "비행단 전장상황도는 상황과 무관하게 항상 1순위로 고정됩니다 — 어떤 CCTV가 지금 떠 있는지 "
          "보여주는 기준 화면이기 때문입니다. 상황 유형에 매핑된 화면이 그다음이고, 남은 자리는 "
          "상시 표출 화면으로 채워 벽에 빈 칸이 남지 않게 합니다."),
         ("역전을 코드로 차단",
          "면적이 우선순위를 따라 단조 감소하도록 배치표를 고정했습니다. 하위 화면이 상위보다 "
          "커지는 일이 생기면 “중요도에 따라 크기를 정한다”는 전제 자체가 무너지기 때문입니다.")],
    )

    d.content(
        "시연", "시연 ② · 작전상황일지 · 상황판",
        "같은 판단으로 기록까지 동시에 완성된다",
        [("작전상황일지 — 사태별 타임라인",
          "14:02  항작       북서방 상공 무인기 2대 식별, 활주로 방향 접근\n"
          "14:03  대공방어   북측 대공포 사격대기 완료, 유도탄 1조 전개\n"
          "14:05  단장       1번 격추 대응, 2번 추적 유지 결심\n\n"
          "CSV로 내보낼 때는 사태당 부서 수만큼 행을 펼쳐 기존 서식에 맞춥니다."),
         ("작전상황판 — 우선순위 현황",
          "#1  긴급   무인기 1번 격추 대응 진행\n"
          "#2  주의   무인기 2번 추적 유지\n"
          "#3  관찰   지상 차량 3대 정상 지원 차량 확인\n\n"
          "복합 상황에서 지휘관·참모의 선택과 집중을 보좌하는 것이 목적입니다.")],
    )

    d.content(
        "시연", "시연 ③ · 전장상황도 자동 마커",
        "언급된 위치에 아이콘이 스스로 놓인다",
        [("자동 배치",
          "발언에 프리셋 키워드(무인기·전술차량·침투 등)가 있으면 해당 아이콘을 자동 배치합니다. "
          "위치는 발언에 언급된 시설명, 없으면 방위를 담당하는 초소로 근사합니다."),
         ("모르면 찍지 않는다",
          "위치를 특정할 수 없으면 아이콘을 올리지 않습니다. 지도에 거짓 정보를 남기지 않기 "
          "위해서입니다. 같은 사태의 대상이 이동하면 새로 쌓지 않고 기존 아이콘을 옮깁니다."),
         ("실무자의 역할",
          "‘전장상황도 조작’ 탭에서 클릭으로 정밀 위치만 보정합니다. 새 마커를 처음부터 손으로 "
          "만들 일은 없습니다.\n\n고정 배치도 격자 A~J × 1~7 · 시설 26곳 · 초소 8곳")],
    )

    d.content(
        "시연", "시연 ④ · 사용자 수동 보정의 이원 반영",
        "같은 실수를 두 번 하지 않는다",
        [("운영 중 — Context Memory에 규칙화 (즉시)",
          "운용자가 잘못된 배치·판단을 교정하면 그 내용이 Context Memory에 즉시 기입되어, "
          "같은 세션과 이후 세션에서 최우선으로 반영됩니다.\n\n"
          "예: “기상 관측 화면은 항상 띄워 두십시오” → 이후 모든 판단에 규칙으로 주입"),
         ("유지보수 시 — 재학습으로 모델에 반영 (주기적)",
          "누적된 보정 사례를 모아 파인튜닝 데이터로 편입하고 주기적으로 재학습해 모델 자체에 "
          "반영합니다. 학습 데이터에도 이 형태를 미리 포함시켰습니다 — “항상 띄워야 하는 화면” "
          "같은 배치 제약조건이 담긴 복합 명령문 6종.")],
        note="기획서 4-⑤ 항의 이원 반영 구조를 그대로 구현했습니다.",
    )

    d.content(
        "파인튜닝", "파인튜닝 · 왜 필요한가",
        "작전 정보는 부대 밖으로 나갈 수 없다",
        [("보안",
          "외부 인터넷이 필요한 상용 AI 대신 부대 내 서버에서 구동하는 오픈소스 sLLM을 써야 "
          "작전 정보 유출 경로가 원천 차단됩니다."),
         ("도메인",
          "범용 모델은 “사태”, “MSEL”, “최소운영활주로” 같은 용어와 상황일지 서식을 모릅니다. "
          "긴 지시문으로 매번 설명하는 대신 모델에 새겨 넣습니다."),
         ("속도",
          "지시문이 짧아지면 입력 토큰이 줄고 그만큼 표출이 빨라집니다. "
          "FAST 20.1% · FULL 40.4% 절감을 실측했습니다.")],
        note="기획서 3-다 요구: 4bit 양자화 2~4GB · 기본 구성 GPU 10~12GB   /   "
             "채택 모델: Qwen2.5-3B-Instruct nf4 ≈ 2.0GB",
    )

    d.table_slide(
        "파인튜닝", "파인튜닝 ① · 데이터 구축", "발언 한 건마다 네 가지 정답을 붙였다",
        ["분할", "시나리오", "발언(턴)", "학습 샘플", "상황 유형"],
        [["train", "176", "1,215", "2,430", "11 / 11"],
         ["valid", "22", "133", "266", "11 / 11"],
         ["test", "22", "179", "358", "11 / 11"],
         ["합계", "220", "1,527", "3,054", ""]],
        col_widths=[0.18, 0.18, 0.20, 0.22, 0.22], highlight_col=2,
        note="기획서 1단계 목표 “학습 데이터 500건 이상” 대비 3배.  라벨 네 종 — 발화자 구분·영향력 태그 / "
             "주요사태(MSEL) 분류 / 사태별 중요도·긴급도 / 화면 구성(COP) 정답.\n"
             "학습 데이터는 앱이 실제로 쓰는 프롬프트 생성 함수를 그대로 호출해 만들므로, 학습과 서빙이 어긋날 수 없습니다.",
    )

    d.content(
        "파인튜닝", "파인튜닝 ② · 학습 구성",
        "무료 GPU 두 장으로 33분에 끝난다",
        [("손실은 응답 구간에만 겁니다",
          "시스템 프롬프트가 길고 모든 샘플에서 동일합니다. 거기에도 손실을 걸면 모델이 프롬프트 "
          "암기에 용량을 쓰고 정작 JSON 형식 준수는 덜 배웁니다.\n\n"
          "손실이 걸리는 토큰 — FAST 3.3% · FULL 8.4%"),
         ("few-shot 없이 학습합니다",
          "파인튜닝의 목적이 긴 예시 없이도 형식을 지키게 만드는 것이므로 학습 데이터에 예시를 "
          "넣지 않습니다. 그래야 서빙에서도 뺄 수 있고 그만큼 빨라집니다.\n\n"
          "입력 토큰 절감 — FAST 20.1% · FULL 40.4%")],
        figures=[("QLoRA", "4bit nf4 · r=16"), ("18.5M", "학습 파라미터 (전체의 1.2%)"),
                 ("33분", "1에폭 152스텝 · T4 ×2 DDP"), ("0.084", "평가 손실 (0.093 → 0.084)")],
    )

    d.metric_slide(
        "성과", "파인튜닝 ③ · 성과", "모든 지표가 동시에 올랐다", METRICS,
        note="일지 누락률 0% (목표 5% 미만) · JSON 유효율 100%.  "
             "측정 조건 — Qwen2.5-1.5B-Instruct, 1에폭, test 40턴, 튜닝 전후 모두 4bit 로드. "
             "튜닝 전에는 few-shot을 켜고 튜닝 후에는 끄고 측정했습니다 — 예시 없이도 형식을 지키는 것 자체가 성과이기 때문입니다.",
    )

    d.table_slide(
        "성과", "기획서 3-라 · 4개 품질 지표 대비 현황", "목표와 현재 위치",
        ["지표", "현행(수동)", "목표", "현재", "상태"],
        [["① 필요정보 표출 지연시간", "약 1~2분", "5초 이내", "4.77초", "달성"],
         ["② 핵심정보 상위배치 정확도", "담당자별 상이", "90% 이상", "88.33%", "근접"],
         ["③ 숙련도별 품질 편차", "편차 큼", "80% 감소", "구조적 0", "설계로 보장"],
         ["④ 작전일지 정확도", "지연·누락 발생", "90% 이상", "84.93%", "근접"],
         ["④ 작전일지 누락률", "지연·누락 발생", "5% 미만", "0%", "달성"]],
        col_widths=[0.32, 0.17, 0.15, 0.14, 0.22], highlight_col=3,
        note="지표 ③이 “구조적 0”인 이유 — 화면 배치를 사람이 아니라 플레이북이 결정하므로, 누가 운용하든 "
             "같은 상황에는 같은 배치가 나옵니다. 숙련도에 따른 편차가 발생할 여지 자체가 없습니다.\n"
             "②와 ④는 목표에 근접했으나 아직 미달입니다. 1.5B 모델을 1에폭 학습한 결과이며, 기획서가 지정한 3B로 올리면 개선 여지가 있습니다.",
    )

    d.flow_slide(
        "성과", "폐쇄망 전환", "코드를 고치지 않고 부대 서버로 옮긴다",
        [("시연 환경", "클라우드 무료 모델 · 인터넷 필요", False),
         ("학습한 모델 병합", "LoRA 어댑터를 베이스에 합쳐 통짜 가중치로", False),
         ("4bit 양자화", "3B Q4_K_M ≈ 2GB · 기획서 요구 부합", False),
         ("부대 내 서버", "vLLM / Ollama · 인터넷 불필요", True)],
        blocks=[("설계 시점부터 대비했습니다",
                 "세 공급자(클라우드 2종, 온프레미스)가 모두 같은 OpenAI 호환 규격을 씁니다. "
                 "사이드바에서 공급자만 바꾸면 접속 주소만 갈아끼워지고 코드 변경이 없습니다."),
                ("시연도 작은 모델로 합니다",
                 "거대 클라우드 모델로 시연해 놓고 온프레미스 12GB에서 된다고 주장하면 성립하지 "
                 "않습니다. 시연에 쓰는 모델도 폐쇄망에서 돌릴 수 있는 크기입니다.")],
    )

    d.content(
        "보안", "보안 · 군 환경 설계",
        "작전 정보를 다루는 AI의 3중 방어",
        [("방어 ① 유출 차단 — 데이터가 부대 밖으로 나가지 않는다",
          "외부 인터넷이 필요한 상용 AI를 쓰지 않습니다. 부대 내 서버에서 구동하는 "
          "오픈소스 sLLM을 직접 학습시켜 쓰므로, 작전 정보가 외부로 전송되는 "
          "경로 자체가 없습니다." + chr(10)*2 + "Qwen2.5-3B · 4bit 약 2GB · 부대 서버 1대로 충분"),
         ("방어 ② 생성 통제 — 모델이 지어낼 수 있는 것이 없다",
          "화면 이름·좌표·제어 명령을 모델이 자유롭게 생성하지 못합니다. LLM은 상황 유형 하나만 "
          "분류하고, 실제 배치는 운용자 플레이북과 폐쇄 카탈로그(270건)를 거칩니다.\n\n"
          "존재하지 않는 화면을 띄우라는 명령이 나올 수 없습니다."),
         ("방어 ③ 권한 한정 — 결심은 사람이 한다",
          "본 체계는 표출·기록 보조에 한정됩니다. 무기체계 통제·사격 지시 등 결심·타격 기능은 "
          "개발 범위에서 명시적으로 제외했습니다. AI는 지휘관이 더 잘 보도록 돕고, 판단은 "
          "지휘관이 합니다.")],
        note="시연 환경에도 접근 통제를 적용했습니다 — 외부 접속에는 암호를 요구하고, 암호가 "
             "설정되지 않았으면 외부 접속 자체를 차단합니다(fail-closed). "
             "API 키는 코드에 넣지 않고 별도 비밀 저장소로 분리했습니다.",
    )

    d.table_slide(
        "보안", "보안 · 예상되는 위험과 대비책", "AI가 지휘소에서 틀렸을 때",
        ["위험 요소", "기획서 대비책", "프로토타입 구현 현황", "상태"],
        [["거짓 정보 생성 (환각)",
          "비정상 수치·맥락에 맞지 않는 화면이 기준치를 초과하면 검증 단계에서 탐지 후 재생성",
          "폐쇄 카탈로그·플레이북으로 생성 경로 자체를 제거. JSON 강제 + 파싱 실패 시 재시도",
          "구현 (유효율 100%)"],
         ["성능 저하 (모델 Drift)",
          "월간 성능 모니터링·경보, 최신 일지·라벨을 추가한 분기별 재학습",
          "평가 하네스 자동화 완료. 데이터 추가 후 같은 명령으로 재측정·재학습",
          "구현"],
         ["실행 불가능한 코드",
          "실행 전 정적 검사 + 테스트 통과를 요구하는 재생성 루프",
          "프로토타입은 코드를 생성하지 않는다. 화면 배치를 코드 생성이 아니라 플레이북 대입으로 해결",
          "해당 없음"],
         ["비정상 시스템 동작",
          "보안 취약점 검사(bandit) + 격리 환경(Docker·Seccomp)에서 허용된 동작만 실행",
          "실제 GUI 자동화(pywinauto) 연동 단계에서 필요. 현재는 화면 조작을 시뮬레이션만 하므로 미적용",
          "향후"]],
        col_widths=[0.19, 0.30, 0.34, 0.17], highlight_col=3,
        note="세 번째 항목이 설계의 핵심입니다. 기획서는 “AI가 만든 제어 코드를 검증하고 격리 실행한다”고 "
             "했지만, 프로토타입은 그 코드를 아예 만들지 않는 구조를 택했습니다 — 위험을 가두는 대신 "
             "위험이 생길 자리를 없앤 것입니다. 실제 Video Wall을 물리적으로 조작하는 단계에 가면 "
             "기획서 원안의 격리 장치가 그때 필요합니다.",
    )

    d.content(
        "마무리", "기대 효과 · 확장성",
        "지휘 결심 속도와 인력 운용",
        [("작전적 효과",
          "보고 직후 지휘관이 눈으로 확인하며 결심합니다. 화면 조작·자료 검색·일지 기입이 "
          "동시에 자동 완성됩니다.\n\n소요시간 1~2분 → 3~5초"),
         ("인력 · 운용 효율",
          "수작업이 줄어 작전병·정통대대 인원의 부담과 오류가 경감되고 핵심 임무로 재배치됩니다.\n\n"
          "부대당 1명 절감 시 연 약 3.6억원 (비행단급 14개 부대 기준)"),
         ("확장성",
          "부대별 최적화는 그 부대의 작전일지로 재학습하는 것만으로 완료됩니다. 구조 자체가 "
          "도메인 데이터만 바꾸면 이식 가능해 육·해군 지휘소, 재난상황실, 소방·경찰 관제센터로 "
          "확산할 수 있습니다.")],
    )

    d.closing(
        "화면 구성 품질이 담당자 숙련도에 의존하지 않게 만드는 일",
        [("4.77초", "발언 종료 → 화면 표출\n목표 5초 이내 달성"),
         ("+30%p", "파인튜닝으로 오른\n상황 분류 정확도"),
         ("2GB", "부대 서버 한 대에\n들어가는 모델 크기")],
        "기획서가 제시한 1단계 “데이터셋 구축과 폐쇄망용 sLLM 파인튜닝”은 완료했습니다. "
        "학습·평가 파이프라인이 자동화되어 있어, 실제 회의록이 확보되면 같은 절차로 "
        "부대에 특화할 수 있습니다.\n\nVOICE-CUE · 팀 작비스 — 감사합니다",
    )

    return d


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path, default=Path("voicecue_발표.pptx"))
    args = ap.parse_args()

    deck = build()
    deck.save(args.out)
    print(f"슬라이드 {deck.n}장 → {args.out} ({args.out.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
